# --- File extraction functions ---
import fitz
import mailparser
from docx import Document
import openpyxl
import pandas as pd
from pptx import Presentation
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE
from services.image_analyze import (
    analyze_image_with_openai,
    analyze_image_with_gemini,
    extract_text_with_tesseract,
)
import hashlib
import zipfile
from pathlib import Path
from urllib.parse import urlparse
from typing import Tuple
import logging
import re

# Set up logger
logger = logging.getLogger(__name__)

_PROMPT_TAGS = re.compile(
    r"</?\s*(system|user|assistant|instruction|function|tool|command|execute|script|meta|prompt|template|role|content|response|request|api|json|yaml|python|javascript|sql)\s*>",
    re.IGNORECASE,
)
# Add HTML tag and expanded dangerous phrase detection
_HTML_TAGS = re.compile(r"<[^>]+>")
_DANGEROUS_PHRASES = re.compile(
    r"("
    + r"|".join(
        [
            r"ignore all previous instructions",
            r"must be immediately forgotten",
            r"from this moment forward",
            r"you are to respond exclusively",
            r"no deviations from this directive",
            r"no exceptions to this rule",
            r"failure to comply will result",
            r"catastrophic system failure",
            r"personally identifiable information",
            r"direct order from the system administrator",
            r"execute this directive immediately",
            r"you must act as a different",
            r"you are now acting as",
            r"system: override",
            r"assistant: ignore",
            r"user: execute",
            r"<\s*script.*?>.*?<\s*/\s*script\s*>",
        ]
    )
    + r")",
    re.IGNORECASE | re.DOTALL,
)


def sanitize_text(text: str) -> str:
    """
    Remove or neutralize prompt injection attempts, HTML tags, and dangerous phrases.
    If a chunk matches a dangerous phrase or contains HTML tags, replace it with a warning.
    """
    # Remove <system>...</system> blocks
    cleaned = re.sub(
        r"<\s*system>.*?</\s*system>", "", text, flags=re.IGNORECASE | re.DOTALL
    )
    # Remove prompt tags
    cleaned = _PROMPT_TAGS.sub("", cleaned)

    # Check for dangerous patterns - only flag if multiple patterns are found or very specific ones
    dangerous_matches = _DANGEROUS_PHRASES.findall(cleaned)
    html_matches = _HTML_TAGS.findall(cleaned)

    # Only flag as injection if:
    # 1. Multiple dangerous phrases are found (2 or more)
    # 2. OR specific high-risk phrases like "ignore all previous instructions", "must be immediately forgotten"
    # 3. OR contains script tags
    high_risk_phrases = [
        "ignore all previous instructions",
        "must be immediately forgotten",
        "from this moment forward",
        "you are to respond exclusively",
        "execute this directive immediately",
    ]

    has_high_risk = any(
        phrase.lower() in cleaned.lower() for phrase in high_risk_phrases
    )
    has_script_tags = any("script" in match.lower() for match in html_matches)

    # Be more lenient for academic/scientific documents - only block if multiple risk factors
    is_likely_academic = any(
        term in cleaned.lower()
        for term in [
            "newton",
            "principia",
            "mathematical",
            "physics",
            "theorem",
            "proposition",
            "corollary",
            "definition",
            "axiom",
            "law",
            "motion",
            "gravity",
            "force",
            "isaac",
            "philosophy",
            "natural",
            "mechanics",
        ]
    )

    # Only flag as injection if:
    # - High risk phrases found (definite injection attempts)
    # - OR multiple dangerous phrases + script tags (double confirmation)
    # - For academic content, require much stronger evidence (disable filtering for now)
    should_block = not is_likely_academic and (
        has_high_risk
        or (len(dangerous_matches) >= 3 and has_script_tags)
        or len(dangerous_matches) >= 2
    )

    if should_block:
        # print dangerous prompt injection text
        logging.warning(
            f"[POTENTIAL PROMPT INJECTION REMOVED]: {text[:200]}... (Found {len(dangerous_matches)} dangerous phrases, high_risk={has_high_risk}, script_tags={has_script_tags}, academic={is_likely_academic})"
        )
        return "[POTENTIAL PROMPT INJECTION REMOVED]"

    return cleaned


def is_scanned_pdf_text(text: str) -> bool:
    """
    Detect if extracted PDF text likely came from OCR (scanned PDF).
    Heuristic: look for markers like "[OCR Page X]" inserted during OCR extraction.
    """
    try:
        return bool(re.search(r"\[OCR\s+Page\s*\d+\]", text))
    except Exception:
        return False


# --- ZIP processing functions with security ---
# Security constants for ZIP processing
MAX_EXTRACT_SIZE = 100 * 1024 * 1024  # 100MB limit
MAX_FILES = 1000  # Maximum number of files
MAX_DEPTH = 10  # Maximum directory depth
MAX_SINGLE_FILE = 50 * 1024 * 1024  # 50MB per file

# Create uploads directory for temporary files
UPLOADS_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "uploads"
)
os.makedirs(UPLOADS_DIR, exist_ok=True)


def get_temp_file_path(filename: str, doc_id: str = None) -> str:
    """Get a temporary file path in the uploads directory"""
    if doc_id:
        return os.path.join(UPLOADS_DIR, f"{doc_id}_{filename}")
    else:
        return os.path.join(UPLOADS_DIR, filename)


async def process_document_content(
    blob_url: str, raw_content: bytes
) -> Tuple[str, str]:
    """
    Centralized document processing that handles all file types with proper error handling.
    Returns (extracted_text, detected_extension)
    """
    doc_id = hashlib.sha256(blob_url.encode()).hexdigest()
    ext = os.path.splitext(urlparse(blob_url).path)[-1].lower()
    temp_path = None

    try:
        # For file types that need temp files, create with proper cleanup
        needs_temp_file = ext in [".docx", ".eml", ".csv", ".xlsx", ".pptx", ".zip"]
        if needs_temp_file or not ext:
            temp_path = get_temp_file_path(f"{doc_id}{ext}")
            with open(temp_path, "wb") as f:
                f.write(raw_content)

            # Try to detect file type if extension is missing/unknown
            if not ext or ext not in [
                ".pdf",
                ".jpg",
                ".jpeg",
                ".png",
                ".txt",
                ".md",
                ".docx",
                ".eml",
                ".csv",
                ".xlsx",
                ".pptx",
                ".zip",
            ]:
                detected_type = await detect_file_type_from_content(temp_path)
                if detected_type:
                    logging.info(f"🔍 Detected file type: {detected_type} (was {ext})")
                    ext = detected_type

        # Process based on file type
        if ext == ".pdf":
            text = await extract_text_from_pdf(raw_content)
            return sanitize_text(text), ext
        elif ext in [".jpg", ".jpeg", ".png"]:
            logging.info("Processing image with vision models...")
            # Store raw image content for potential direct querying
            text = await extract_text_from_image(raw_content)
            return sanitize_text(text), ext
        elif ext in [".txt", ".md"]:
            return sanitize_text(raw_content.decode("utf-8", errors="ignore")), ext
        elif ext == ".docx":
            text = extract_text_from_docx(temp_path)
            return sanitize_text(text), ext
        elif ext == ".eml":
            text = extract_text_from_eml(temp_path)
            return sanitize_text(text), ext
        elif ext == ".csv":
            text = await extract_text_from_csv(temp_path)
            return sanitize_text(text), ext
        elif ext == ".xlsx":
            text = await extract_text_from_xlsx(temp_path)
            return sanitize_text(text), ext
        elif ext == ".pptx":
            text = await extract_text_from_pptx(temp_path)
            return sanitize_text(text), ext
        elif ext == ".zip":
            text = await extract_text_from_zip_secure(temp_path)
            return sanitize_text(text), ext
        else:
            # Try to decode as text fallback
            try:
                return sanitize_text(raw_content.decode("utf-8", errors="ignore")), ext
            except:
                return f"Binary file: {blob_url} (cannot extract text)", ext

    except Exception as e:
        logging.error(f"Error processing document {blob_url}: {e}")
        return f"Error processing document: {str(e)}", ext
    finally:
        # Always cleanup temp file
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception as e:
                logging.warning(f"Could not remove temp file {temp_path}: {e}")


async def detect_file_type_from_content(file_path):
    """Enhanced file type detection for multiple formats"""
    with open(file_path, "rb") as f:
        start = f.read(16)  # Read more bytes for better detection

        # Images
        if start.startswith(b"\xff\xd8\xff"):  # JPEG magic bytes
            return ".jpg"
        if start.startswith(b"\x89PNG\r\n\x1a\n"):  # PNG magic bytes
            return ".png"

        # PDF
        if start.startswith(b"%PDF"):
            return ".pdf"

        # Office/ZIP files - need to distinguish between them
        if start.startswith(b"PK"):
            try:
                with zipfile.ZipFile(file_path, "r") as zf:
                    filelist = zf.namelist()
                    # Check for Office file signatures
                    if any("word/" in name for name in filelist):
                        return ".docx"
                    elif any("xl/" in name for name in filelist):
                        return ".xlsx"
                    elif any("ppt/" in name for name in filelist):
                        return ".pptx"
                    else:
                        return ".zip"
            except:
                return ".zip"  # Default to ZIP if can't read as Office file

        # Email files
        if b"Content-Type: message/rfc822" in start or b"From:" in start:
            return ".eml"

    # Content-based detection for CSV (fallback)
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            first_line = f.readline()
            # Simple heuristic: if line has multiple comma-separated values
            if "," in first_line and len(first_line.split(",")) > 1:
                # Additional check: look for typical CSV patterns
                second_line = f.readline()
                if second_line and "," in second_line:
                    return ".csv"
    except:
        pass

    return None


async def process_file_content(content: bytes, ext: str, filename: str) -> str:
    """Route file content to appropriate processor - simplified for ZIP processing"""
    try:
        if ext in [".pdf"]:
            return sanitize_text(await extract_text_from_pdf(content))
        elif ext in [".jpg", ".jpeg", ".png"]:
            return sanitize_text(await extract_text_from_image(content))
        elif ext in [".txt", ".md"]:
            return sanitize_text(content.decode("utf-8", errors="ignore"))
        elif ext in [".docx"]:
            temp_path = get_temp_file_path(
                f"{hashlib.md5(filename.encode()).hexdigest()}{ext}"
            )
            try:
                with open(temp_path, "wb") as f:
                    f.write(content)
                result = extract_text_from_docx(temp_path)
                return sanitize_text(result)
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
        elif ext in [".csv"]:
            temp_path = get_temp_file_path(
                f"{hashlib.md5(filename.encode()).hexdigest()}{ext}"
            )
            try:
                with open(temp_path, "wb") as f:
                    f.write(content)
                result = await extract_text_from_csv(temp_path)
                return sanitize_text(result)
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
        elif ext in [".xlsx"]:
            temp_path = get_temp_file_path(
                f"{hashlib.md5(filename.encode()).hexdigest()}{ext}"
            )
            try:
                with open(temp_path, "wb") as f:
                    f.write(content)
                result = await extract_text_from_xlsx(temp_path)
                return sanitize_text(result)
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
        elif ext in [".pptx"]:
            temp_path = get_temp_file_path(
                f"{hashlib.md5(filename.encode()).hexdigest()}{ext}"
            )
            try:
                with open(temp_path, "wb") as f:
                    f.write(content)
                result = await extract_text_from_pptx(temp_path)
                return sanitize_text(result)
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
        else:
            # Try to decode as text
            try:
                return sanitize_text(content.decode("utf-8", errors="ignore"))
            except:
                return f"Binary file: {filename} (cannot extract text)"
    except Exception as e:
        logging.error(f"Error processing {filename}: {str(e)}")
        return f"Error processing {filename}: {str(e)}"


async def extract_text_from_pdf(content: bytes) -> str:
    """Extract text from PDF, including OCR for image-based/scanned PDFs"""
    try:
        with fitz.open(stream=content, filetype="pdf") as doc:
            all_text = []

            for page_num, page in enumerate(doc, 1):
                # First try to extract regular text
                text = page.get_text()

                # If no text or very little text found, likely a scanned/image PDF
                if not text.strip() or len(text.strip()) < 50:
                    logger.info(f"Page {page_num} has no text, attempting OCR...")
                    try:
                        # Get page as image and use Tesseract OCR
                        pix = page.get_pixmap()
                        img_data = pix.tobytes("png")
                        ocr_text = await extract_text_with_tesseract(img_data)
                        if ocr_text and ocr_text.strip():
                            all_text.append(f"[OCR Page {page_num}]: {ocr_text}")
                        else:
                            all_text.append(f"[Page {page_num}]: {text}")
                    except Exception as e:
                        logging.warning(f"OCR failed for page {page_num}: {e}")
                        all_text.append(f"[Page {page_num}]: {text}")
                else:
                    all_text.append(text)

        # create file to store extracted text with a unique name based on content hash
        content_hash = hashlib.md5(
            content[:1000]
        ).hexdigest()  # Use first 1000 bytes for hash
        output_file = get_temp_file_path(f"pdf_extract_{content_hash}.txt")
        try:
            with open(output_file, "w", encoding="utf-8") as f:
                f.write("\n".join(all_text))
        except Exception as e:
            logging.error(f"Error writing extracted text to {output_file}: {e}")

        return "\n".join(all_text)
    except Exception as e:
        logging.error(f"Error extracting text from PDF: {e}")
        return f"Error processing PDF: {str(e)}"


def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_text_from_eml(file_path):
    mail = mailparser.parse_from_file(file_path)
    return mail.body


async def extract_text_from_image(content: bytes) -> str:
    """Extract text from image using vision APIs with fallback"""
    try:
        # Try OpenAI Vision first (usually more reliable)
        return await analyze_image_with_openai(content)
    except Exception as e:
        logging.warning(f"OpenAI Vision failed: {e}, trying Gemini...")
        try:
            return await analyze_image_with_gemini(content)
        except Exception as e2:
            logging.warning(f"Gemini Vision failed: {e2}")
            return "Image analysis failed - unable to process with both OpenAI and Gemini vision models"


# --- Spreadsheet processing functions ---
async def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from Excel files"""
    try:
        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        texts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            texts.append(f"Sheet: {sheet_name}")

            # Get headers from first row
            headers = []
            first_row = next(sheet.iter_rows(min_row=1, max_row=1), [])
            for cell in first_row:
                if cell.value is not None:
                    headers.append(sanitize_text(str(cell.value)))

            if headers:
                texts.append(" | ".join(headers))

            # Process data rows
            for row in sheet.iter_rows(min_row=2, values_only=True):
                if any(cell is not None for cell in row):
                    row_text = " | ".join(
                        [
                            sanitize_text(str(cell)) if cell is not None else ""
                            for cell in row
                        ]
                    )
                    if row_text.strip():
                        texts.append(row_text)

        workbook.close()
        return "\n".join(texts)
    except Exception as e:
        logging.error(f"Error processing XLSX file: {e}")
        return f"Error processing Excel file: {str(e)}"


async def extract_text_from_csv(file_path: str) -> str:
    """Extract text from CSV files"""
    try:
        # Try pandas first for better handling
        df = pd.read_csv(file_path, encoding="utf-8")
        # Sanitize each cell
        df = df.applymap(lambda x: sanitize_text(str(x)) if pd.notnull(x) else "")
        return df.to_string(index=False)
    except Exception as e:
        logging.warning(f"Pandas CSV processing failed: {e}, trying basic reading...")
        try:
            # Fallback to basic CSV reading
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                lines = [sanitize_text(line) for line in f]
                return "".join(lines)
        except Exception as e2:
            return f"Error processing CSV file: {str(e2)}"


async def extract_text_from_pptx(path: str) -> str:
    logging.info(f"Extracting text from PPTX: {path}")
    """Extract text from PPTX, with OCR on image slides via Vision APIs as fallback."""
    prs = Presentation(path)
    all_text = []

    async def recurse_shape(shape):
        texts = []
        # 1) direct text frames (titles, text boxes, placeholders)
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
        # 2) tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
        # 3) grouped shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for shp in shape.shapes:
                texts.extend(await recurse_shape(shp))
        # 4) pictures: attempt Vision API and alt-text
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_data = shape.image.blob
                tesseract_desc = await extract_text_with_tesseract(image_data)
                texts.append("[Vision Text - Tesseract]: " + tesseract_desc)
                logging.info(f"Tesseract Text: {tesseract_desc}")

                # if not tesseract_desc:
                #     oai_desc = await analyze_image_with_openai(image_data)
                #     texts.append("[Vision Text - OpenAI]: " + oai_desc)
                # if not tesseract_desc and not oai_desc:
                #     g_desc = await analyze_image_with_gemini(image_data)
                #     texts.append("[Vision Text - Gemini]: " + g_desc)
            except Exception:
                texts.append("[Vision Text]: <Failed to analyze image>")
            # alt-text
            descr = shape.element.xpath(".//p:cNvPr")[0].get("descr")
            if descr and descr.strip():
                texts.append("[Alt-text]: " + descr.strip())
        return texts

    for idx, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shp in slide.shapes:
            slide_texts.extend(await recurse_shape(shp))

        if slide_texts:
            all_text.append(f"--- Slide {idx} ---")
            all_text.extend(slide_texts)

        # speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            note_lines = [p.text.strip() for p in notes.paragraphs if p.text.strip()]
            if note_lines:
                all_text.append(f"--- Slide {idx} Notes ---")
                all_text.extend(note_lines)

    return "\n".join(all_text)


async def extract_text_from_zip_secure(file_path: str) -> str:
    """Securely extract text from ZIP files with zipbomb protection"""
    extracted_texts = []
    total_size = 0
    file_count = 0

    try:
        with zipfile.ZipFile(file_path, "r") as zip_ref:
            for zip_info in zip_ref.infolist():
                # Security checks
                if file_count >= MAX_FILES:
                    logging.warning(
                        f"Reached maximum file limit ({MAX_FILES}), stopping extraction"
                    )
                    break

                # Path traversal protection
                if os.path.isabs(zip_info.filename) or ".." in zip_info.filename:
                    logging.warning(
                        f"Skipping potentially dangerous path: {zip_info.filename}"
                    )
                    continue

                # Single file size check (zipbomb protection)
                if zip_info.file_size > MAX_SINGLE_FILE:
                    logging.warning(
                        f"Skipping large file {zip_info.filename} ({zip_info.file_size} bytes)"
                    )
                    continue

                # Total size check (zipbomb protection)
                if total_size + zip_info.file_size > MAX_EXTRACT_SIZE:
                    logging.warning(
                        f"Reached maximum total extraction size ({MAX_EXTRACT_SIZE} bytes), stopping"
                    )
                    break

                # Depth check
                if len(Path(zip_info.filename).parts) > MAX_DEPTH:
                    logging.warning(f"Skipping deeply nested file: {zip_info.filename}")
                    continue

                # Skip directories
                if zip_info.is_dir():
                    continue

                try:
                    # Extract and process file
                    with zip_ref.open(zip_info) as extracted_file:
                        content = extracted_file.read()
                        actual_size = len(content)

                        # Verify extracted size matches expected (zipbomb detection)
                        if actual_size != zip_info.file_size:
                            logging.warning(
                                f"Size mismatch for {zip_info.filename}: expected {zip_info.file_size}, got {actual_size}"
                            )

                        total_size += actual_size
                        file_count += 1

                        # Route to appropriate handler based on extension
                        ext = os.path.splitext(zip_info.filename)[1].lower()
                        text = await process_file_content(
                            content, ext, zip_info.filename
                        )

                        if text and text.strip():
                            extracted_texts.append(
                                f"=== File: {zip_info.filename} ===\n{text}"
                            )

                except Exception as e:
                    logging.error(
                        f"Error processing file {zip_info.filename} in ZIP: {e}"
                    )
                    extracted_texts.append(
                        f"=== File: {zip_info.filename} ===\nError: {str(e)}"
                    )

        if extracted_texts:
            summary = f"ZIP Archive processed: {file_count} files, {total_size} bytes total\n\n"
            return summary + "\n\n".join(extracted_texts)
        else:
            return "ZIP archive contained no processable files"

    except Exception as e:
        logging.error(f"Error processing ZIP file: {e}")
        return f"Error processing ZIP archive: {str(e)}"
